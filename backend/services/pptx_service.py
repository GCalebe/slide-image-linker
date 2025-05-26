
from pathlib import Path
from typing import List, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

EMU_PER_PX = 9525  # 1 pixel ≈ 9525 EMU

def _shape_bbox(shp) -> Tuple[int, int, int, int]:
    return shp.left // EMU_PER_PX, shp.top // EMU_PER_PX, shp.width // EMU_PER_PX, shp.height // EMU_PER_PX

def extract_slide_png(pptx_path: Path, slide_number: int) -> Tuple[bytes, List[dict]]:
    """Gera PNG dummy do slide e lista de bounding-boxes dos shapes relevantes."""
    prs = Presentation(pptx_path)
    if slide_number < 1 or slide_number > len(prs.slides):
        slide_number = 1
    slide = prs.slides[slide_number - 1]

    width_px = prs.slide_width // EMU_PER_PX
    height_px = prs.slide_height // EMU_PER_PX

    img = Image.new("RGB", (width_px, height_px), (240, 240, 240))

    meta = []
    for shp in slide.shapes:
        if shp.shape_type in (
            MSO_SHAPE_TYPE.TEXT_BOX,
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.PLACEHOLDER,
        ):
            x, y, w, h = _shape_bbox(shp)
            meta.append(
                dict(
                    id=f"Slide{slide_number}-{shp.shape_id}",
                    x=x,
                    y=y,
                    w=w,
                    h=h,
                )
            )

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue(), meta

def apply_mappings_to_pptx(pptx_path: Path, mappings: List[dict]) -> bytes:
    """Insere textos nos shapes conforme mappings."""
    prs = Presentation(pptx_path)

    shape_lookup = {
        f"Slide{slide_idx+1}-{shape.shape_id}": shape
        for slide_idx, slide in enumerate(prs.slides)
        for shape in slide.shapes
    }

    for mp in mappings:
        sid = mp.get("shape_id")
        new_text = mp.get("new_text")
        if sid in shape_lookup and new_text is not None:
            shape = shape_lookup[sid]
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                shape.text_frame.text = str(new_text)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
